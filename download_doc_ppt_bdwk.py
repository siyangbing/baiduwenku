import os
import time

from selenium import webdriver
from selenium.webdriver.common.desired_capabilities import DesiredCapabilities
from scrapy import Selector
import requests
from my_fake_useragent import UserAgent
import docx
from docx.shared import Inches
import cv2
from pptx import Presentation
from pptx.util import Inches

#dows是的chromedriver
# chromedriver_path = "./chromedriver.exe"
#用ubuntu的chromedriver
chromedriver_path = "./chromedriver"

doc_dir_path = "./doc"
ppt_dir_path = "./ppt"
# url = "https://wenku.baidu.com/view/4410199cb0717fd5370cdc2e.html?fr=search"# doc_txt p
# url = "https://wenku.baidu.com/view/4d18916f7c21af45b307e87101f69e314332fa36.html" # doc_txt span
# url = "https://wenku.baidu.com/view/dea519c7e53a580216fcfefa.html?fr=search" # doc_txt span br
# url = 'https://wk.baidu.com/view/062edabeb6360b4c2e3f5727a5e9856a5712262d?pcf=2&bfetype=new' # doc_img
# url = "https://wenku.baidu.com/view/2af6de34a7e9856a561252d380eb6294dd88228d"# vip限定doc
# url = "https://wenku.baidu.com/view/3de365cc6aec0975f46527d3240c844769eaa0aa.html?fr=search" #ppt
# url = "https://wenku.baidu.com/view/18a8bc08094e767f5acfa1c7aa00b52acec79c55"#pdf
# url = "https://wenku.baidu.com/view/bbe27bf21b5f312b3169a45177232f60dccce772"
# url = "https://wenku.baidu.com/view/5cb11d096e1aff00bed5b9f3f90f76c660374c24.html?fr=search"
# url = "https://wenku.baidu.com/view/71f9818fef06eff9aef8941ea76e58fafab045a6.html"
# url = "https://wenku.baidu.com/view/ffc6b32a68eae009581b6bd97f1922791788be69.html"
url = "https://wenku.baidu.com/view/ea815d9c1fd9ad51f01dc281e53a580217fc50e0.html"

class DownloadImg():
    def __init__(self):
        self.ua = UserAgent()

    def download_one_img(self, img_url, saved_path):
        # 下载图片
        header = {
            "User-Agent": "{}".format(self.ua.random().strip()),
            'Connection': 'close'}
        r = requests.get(img_url, headers=header, stream=True)
        print("请求图片状态码　{}".format(r.status_code))  # 返回状态码
        if r.status_code == 200:  # 写入图片
            with open(saved_path, mode="wb") as f:
                f.write(r.content)
            print("download {} success!".format(saved_path))
        del r
        return saved_path


class StartChrome():
    def __init__(self):
        mobile_emulation = {"deviceName": "Galaxy S5"}
        capabilities = DesiredCapabilities.CHROME
        capabilities['loggingPrefs'] = {'browser': 'ALL'}
        options = webdriver.ChromeOptions()
        options.add_experimental_option("mobileEmulation", mobile_emulation)
        self.brower = webdriver.Chrome(executable_path=chromedriver_path, desired_capabilities=capabilities,
                                       chrome_options=options)
        # 启动浏览器，打开需要下载的网页
        self.brower.get(url)
        self.download_img = DownloadImg()

    def click_ele(self, click_xpath):
        # 单击指定控件
        click_ele = self.brower.find_elements_by_xpath(click_xpath)
        if click_ele:
            click_ele[0].location_once_scrolled_into_view  # 滚动到控件位置
            self.brower.execute_script('arguments[0].click()', click_ele[0])  # 单击控件，即使控件被遮挡，同样可以单击

    def judge_doc(self, contents):
        # 判断文档类别
        p_list = ''.join(contents.xpath("./text()").extract())
        span_list = ''.join(contents.xpath("./span/text()").extract())
        # # if span_list
        # if len(span_list)>len(p_list):
        #     xpath_content_one = "./br/text()|./span/text()|./text()"
        # elif len(span_list)<len(p_list):
        #     # xpath_content_one = "./br/text()|./text()"
        #     xpath_content_one = "./br/text()|./span/text()|./text()"
        if len(span_list)!=len(p_list):
            xpath_content_one = "./br/text()|./span/text()|./text()"
        else:
            xpath_content_one = "./span/img/@src"
        return xpath_content_one

    def create_ppt_doc(self, ppt_dir_path, doc_dir_path):
        # 点击关闭开通会员按钮
        xpath_close_button = "//div[@class='na-dialog-wrap show']/div/div/div[@class='btn-close']"
        self.click_ele(xpath_close_button)
        # 点击继续阅读
        xpath_continue_read_button = "//div[@class='foldpagewg-icon']"
        self.click_ele(xpath_continue_read_button)
        # 点击取消打开百度ａｐｐ按钮
        xpath_next_content_button = "//div[@class='btn-wrap']/div[@class='btn-cancel']"
        self.click_ele(xpath_next_content_button)
        # 循环点击加载更多按钮，直到显示全文
        click_count = 0
        while True:
            # 如果到了最后一页就跳出循环
            if self.brower.find_elements_by_xpath("//div[@class='pagerwg-loadSucc hide']") or self.brower.find_elements_by_xpath("//div[@class='pagerwg-button' and @style='display: none;']"):
                break
            # 点击加载更多
            xpath_loading_more_button = "//span[@class='pagerwg-arrow-lower']"
            self.click_ele(xpath_loading_more_button)
            click_count += 1
            print("第{}次点击加载更多!".format(click_count))
            # 等待一秒，等浏览器加载
            time.sleep(1.5)

        # 获取html内容
        sel = Selector(text=self.brower.page_source)
        #判断文档类型
        xpath_content = "//div[@class='content singlePage wk-container']/div/p/img/@data-loading-src|//div[@class='content singlePage wk-container']/div/p/img/@data-src"
        contents = sel.xpath(xpath_content).extract()
        if contents:#如果是ppt
            self.create_ppt(ppt_dir_path, sel)
        else:#如果是doc
            self.create_doc(doc_dir_path, sel)
        # a = 3333
        # return sel

    def create_ppt(self, ppt_dir_path, sel):
        # 如果文件夹不存在就创建一个
        if not os.path.exists(ppt_dir_path):
            os.makedirs(ppt_dir_path)

        SLD_LAYOUT_TITLE_AND_CONTENT = 6  # 6代表ppt模版为空
        prs = Presentation()  # 实例化ppt

        # # 获取完整html
        # sel = self.get_html_data()
        # 获取标题
        xpath_title = "//div[@class='doc-title']/text()"
        title = "".join(sel.xpath(xpath_title).extract()).strip()
        # 获取内容
        xpath_content_p = "//div[@class='content singlePage wk-container']/div/p/img"
        xpath_content_p_list = sel.xpath(xpath_content_p)
        xpath_content_p_url_list=[]
        for imgs in xpath_content_p_list:
            xpath_content = "./@data-loading-src|./@data-src|./@src"
            contents_list = imgs.xpath(xpath_content).extract()
            xpath_content_p_url_list.append(contents_list)

        img_path_list = []  # 保存下载的图片路径，方便后续图片插入ppt和删除图片
        # 下载图片到指定目录
        for index, content_img_p in enumerate(xpath_content_p_url_list):
            p_img_path_list=[]
            for index_1,img_one in enumerate(content_img_p):
                one_img_saved_path = os.path.join(ppt_dir_path, "{}_{}.jpg".format(index,index_1))
                self.download_img.download_one_img(img_one, one_img_saved_path)
                p_img_path_list.append(one_img_saved_path)

            p_img_max_shape = 0
            for index,p_img_path in enumerate(p_img_path_list):
                img_shape = cv2.imread(p_img_path).shape
                if p_img_max_shape<img_shape[0]:
                    p_img_max_shape = img_shape[0]
                    index_max_img = index
            img_path_list.append(p_img_path_list[index_max_img])


        print(img_path_list)
        # 获取下载的图片中最大的图片的尺寸
        img_shape_max=[0,0]
        for img_path_one in img_path_list:
            img_path_one_shape = cv2.imread(img_path_one).shape
            if img_path_one_shape[0]>img_shape_max[0]:
                img_shape_max = img_path_one_shape
        # 把图片统一缩放最大的尺寸
        for img_path_one in img_path_list:
            cv2.imwrite(img_path_one,cv2.resize(cv2.imread(img_path_one),(img_shape_max[1],img_shape_max[0])))
        # img_shape_path = img_path_list[0]
        # 获得图片的尺寸
        # img_shape = cv2.imread(img_shape_path).shape
        # 把像素转换为ppt中的长度单位emu,默认dpi是720
        # 1厘米=28.346像素=360000
        # 1像素 = 12700emu
        prs.slide_width = img_shape_max[1] * 12700  # 换算单位
        prs.slide_height = img_shape_max[0] * 12700

        for img_path_one in img_path_list:
            left = Inches(0)
            right = Inches(0)
            # width = Inches(1)
            slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
            slide = prs.slides.add_slide(slide_layout)
            pic = slide.shapes.add_picture(img_path_one, left, right, )
            print("insert {} into pptx success!".format(img_path_one))
            # os.remove(img_path_one)

        for root,dirs,files in os.walk(ppt_dir_path):
            for file in files:
                if file.endswith(".jpg"):
                    img_path = os.path.join(root,file)
                    os.remove(img_path)

        prs.save(os.path.join(ppt_dir_path, title + ".pptx"))
        print("download {} success!".format(os.path.join(ppt_dir_path, title + ".pptx")))

    def create_doc(self, doc_dir_path, sel):
        # 如果文件夹不存在就创建一个
        if not os.path.exists(doc_dir_path):
            os.makedirs(doc_dir_path)
        # # 获取完整html
        # sel = self.get_html_data()
        # 获取标题
        xpath_title = "//div[@class='doc-title']/text()"
        title = "".join(sel.xpath(xpath_title).extract()).strip()

        document = docx.Document()  # 创建word文档
        document.add_heading(title, 0)  # 添加标题

        # 获取文章内容
        xpath_content = "//div[contains(@data-id,'div_class_')]//p"
        # xpath_content = "//div[contains(@data-id,'div_class_')]/p"
        contents = sel.xpath(xpath_content)
        # 判断内容类别
        xpath_content_one = self.judge_doc(contents)
        if xpath_content_one.endswith("text()"):  # 如果是文字就直接爬
            for content_one in contents:
                one_p_list = content_one.xpath(xpath_content_one).extract()
                p_txt = ""
                for p in one_p_list:
                    if p==" ":
                        p_txt += ('\n'+p)
                    else:
                        p_txt += p
                # content_txt_one = '*'.join(content_one.xpath(xpath_content_one).extract())
                pp = document.add_paragraph(p_txt)
            document.save(os.path.join(doc_dir_path, '{}.docx'.format(title)))
            print("download {} success!".format(title))
        elif xpath_content_one.endswith("@src"):  # 如果是图片就下载图片
            for index, content_one in enumerate(contents.xpath(xpath_content_one).extract()):
                # 获取图片下载路径
                content_img_one_url = 'https:' + content_one
                # 保存图片
                saved_image_path = self.download_img.download_one_img(content_img_one_url, os.path.join(doc_dir_path,
                                                                                                        "{}.jpg".format(
                                                                                                            index)))
                document.add_picture(saved_image_path, width=Inches(6))  # 在文档中加入图片
                os.remove(saved_image_path)  # 删除下载的图片
            document.save(os.path.join(doc_dir_path, '{}.docx'.format(title)))  # 保存文档到指定位置
            print("download {} success!".format(title))


if __name__ == "__main__":
    start_chrome = StartChrome()
    # start_chrome.create_doc_txt(doc_dir_path)
    start_chrome.create_ppt_doc(ppt_dir_path, doc_dir_path)
